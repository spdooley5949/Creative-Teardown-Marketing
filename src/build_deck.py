"""Creative Teardown — Friday deck. 10 slides, 16:9, light palette."""
from pptx import Presentation
from pptx.util import Inches as I, Pt
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
import sys

# Light palette. Deep slate rather than black; one warm signal reserved
# for the whitespace finding so the empty column reads as the point.
PAPER = C(0xFF, 0xFF, 0xFF)   # slide ground
WASH  = C(0xEC, 0xF1, 0xF6)   # title / close ground, light but distinct
MIST  = C(0xF4, 0xF6, 0xF8)   # panels
INK   = C(0x1F, 0x29, 0x33)   # primary text
MID   = C(0x44, 0x51, 0x5F)   # secondary
GREY  = C(0x69, 0x74, 0x82)   # tertiary
LINE  = C(0xD5, 0xDD, 0xE4)
SIG   = C(0xC2, 0x41, 0x0C)   # signal: the open lane
TINT  = C(0xFD, 0xEC, 0xE2)   # signal panel
SOFTB = C(0xDF, 0xE7, 0xEF)   # rule on wash

HEAD, BODY = "Cambria", "Calibri"
W, H, M = 13.333, 7.5, 0.7

prs = Presentation()
prs.slide_width, prs.slide_height = I(W), I(H)
BLANK = prs.slide_layouts[6]


def slide(bg=PAPER):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def text(s, t, x, y, w, h, *, size=14, font=BODY, color=INK, bold=False,
         italic=False, align=PP_ALIGN.LEFT, spacing=None):
    box = s.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, ln in enumerate(str(t).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing:
            p.line_spacing = Pt(spacing)
        r = p.add_run()
        r.text = ln
        r.font.size, r.font.name, r.font.bold, r.font.italic = Pt(size), font, bold, italic
        r.font.color.rgb = color
    return box


def rect(s, x, y, w, h, fill):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(x), I(y), I(w), I(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh


def circle(s, x, y, d, fill):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, I(x), I(y), I(d), I(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh


def eyebrow(s, t):
    text(s, t.upper(), M, 0.34, W - M * 2, 0.28, size=11, bold=True, color=SIG)


def title(s, t, size=36, y=0.62):
    text(s, t, M, y, W - M * 2, 1.05, size=size, font=HEAD, bold=True, color=INK)


def notes(s, t):
    s.notes_slide.notes_text_frame.text = t


# ── 1 · TITLE ───────────────────────────────────────────────────────────────
s = slide(WASH)
text(s, "Twelve brands.\nOne message.", M, 1.5, 9.4, 2.6,
     size=54, font=HEAD, bold=True, color=INK, spacing=58)
text(s, "A creative teardown of what fitness apparel actually says in paid search",
     M, 4.2, 9.4, 0.5, size=17, color=MID)
rect(s, M, 5.05, 1.5, 0.045, SIG)
text(s, "101 text ads   ·   12 brands   ·   listed 12 August 2026   ·   replicated on a second 1,786-ad set",
     M, 5.33, 11.4, 0.35, size=13, color=SIG)
text(s, "Shane Dooley", M, 6.2, 6, 0.3, size=14, bold=True, color=INK)
text(s, "with Andrew Silver, Catherine Chambers, Jessica Smith, Jesse Padilla and William Levine",
     M, 6.55, 11.5, 0.3, size=11.5, color=GREY)
notes(s, "Ten minutes. Open on the finding, not an agenda. Twelve fitness apparel brands, "
         "what they are actually saying in Google search.")

# ── 2 · THE FINDING ─────────────────────────────────────────────────────────
s = slide()
eyebrow(s, "The finding")
title(s, "The category has run out of ways to say “comfortable”")
for i, (n, lab, hot) in enumerate([
        ("59%", "claim technical\nperformance", False),
        ("58%", "claim comfort", False),
        ("51%", "claim gym-to-street\nversatility", False),
        ("6%", "claim community\nor belonging", True)]):
    x = M + i * 3.02
    rect(s, x, 2.35, 2.8, 2.55, TINT if hot else MIST)
    text(s, n, x + 0.26, 2.62, 2.3, 1.15, size=52, font=HEAD, bold=True,
         color=SIG if hot else INK)
    text(s, lab, x + 0.26, 3.82, 2.3, 0.9, size=13, color=SIG if hot else MID)
text(s, "Across all twelve brands the language is near-interchangeable: "
        "“buttery soft,” “built to move in,” “looks as good as it feels.”",
     M, 5.3, W - M * 2, 0.45, size=15, color=MID)
text(s, "Community and belonging is six ads in the entire set. Five of them are Gymshark’s.",
     M, 5.85, W - M * 2, 0.4, size=15, bold=True, color=SIG)
notes(s, "Lead here. Three saturated claims, then the cliff. Community is 6 ads, 5 of them Gymshark. Do not explain method yet.")

# ── 3 · METHOD ──────────────────────────────────────────────────────────────
s = slide()
eyebrow(s, "Method")
title(s, "101 text ads, twelve brands, one grid")
for i, (k, v) in enumerate([
        ("Collect", "Every search text ad the Google Ads Transparency Center listed for each "
                    "brand on 12 August, transcribed by hand. Public source, no login, no scraping."),
        ("Classify", "Claude tags each ad; tags are frozen to a file so numbers reproduce "
                     "exactly. A teammate then audited all 101 rows and 42 corrections were applied."),
        ("Assemble", "Counts roll into a matrix as a percent of each brand’s ads, so brands "
                     "sampled at different depths stay comparable.")]):
    y = 2.32 + i * 1.3
    circle(s, M, y, 0.5, INK)
    text(s, str(i + 1), M, y + 0.09, 0.5, 0.35, size=15, bold=True,
         color=PAPER, align=PP_ALIGN.CENTER)
    text(s, k, M + 0.82, y - 0.02, 2.0, 0.35, size=16, bold=True)
    text(s, v, M + 2.6, y - 0.02, 9.25, 1.0, size=13, color=MID)
rect(s, M, 6.28, W - M * 2, 0.72, MIST)
text(s, "lululemon 18  ·  On 12  ·  Arc’teryx 12  ·  Vuori 10  ·  Gymshark 10  ·  Alo Yoga 9  ·  "
        "Nike 8  ·  Under Armour 7  ·  Tracksmith 4  ·  New Balance 4  ·  Fabletics 4  ·  Patagonia 3",
     M + 0.28, 6.45, W - M * 2 - 0.56, 0.4, size=11, color=GREY)
notes(s, "Thirty seconds. If asked: no date filter was applied, so this is every text ad the "
         "Transparency Center listed that day, not a verified set of live ads. Does not change the finding.")

# ── 4 · THE MATRIX ──────────────────────────────────────────────────────────
s = slide()
eyebrow(s, "Evidence")
title(s, "Where the category is stacked, and where it is not")
cd = CategoryChartData()
cd.categories = ["Community", "Sustainability", "Innovation", "Price / Value",
                 "Style", "Versatility", "Comfort", "Performance"]
cd.add_series("Share of ads", (6, 7, 14, 34, 36, 51, 58, 59))
ch = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, I(M), I(2.15),
                        I(W - M * 2), I(4.05), cd).chart
ch.has_legend = False; ch.has_title = False
pl = ch.plots[0]; pl.gap_width = 45; pl.has_data_labels = True
dl = pl.data_labels
dl.number_format, dl.number_format_is_linked = '0"%"', False
dl.position = XL_LABEL_POSITION.OUTSIDE_END
dl.font.size, dl.font.name, dl.font.bold = Pt(12), BODY, True
dl.font.color.rgb = INK
for idx in range(8):
    pt = pl.series[0].points[idx]
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = SIG if idx == 0 else (GREY if idx < 4 else INK)
ca, va = ch.category_axis, ch.value_axis
ca.has_major_gridlines = False
ca.tick_labels.font.size, ca.tick_labels.font.name = Pt(13), BODY
ca.tick_labels.font.color.rgb = INK
ca.format.line.color.rgb = LINE
va.has_major_gridlines = False; va.visible = False
text(s, "Percent of all 101 ads using each claim theme. An ad can carry more than one.",
     M, 6.45, W - M * 2, 0.35, size=11, color=GREY)
notes(s, "The shape is the argument. Four crowded claims, then a cliff. "
         "Community in orange at the bottom is where the deck is heading.")

# ── 5 · REPLICATION ─────────────────────────────────────────────────────────
s = slide()
eyebrow(s, "Independent check")
title(s, "The finding holds on data we did not collect")
text(s, "A teammate contributed 1,786 ads across eight brands, 2021 to 2026, gathered a "
        "different way. Four brands appear in both sets. Neither collector saw the other’s work.",
     M, 2.12, 11.6, 0.7, size=14, color=MID)
hdr = ["Brand", "Ads, by hand", "Ads, harvested", "Gap, all themes", "Community"]
xs = [M, 4.3, 6.5, 8.7, 11.0]
for i, hd in enumerate(hdr):
    text(s, hd.upper(), xs[i], 3.05, 2.1, 0.3, size=10, bold=True,
         color=SIG if i == 4 else GREY)
rect(s, M, 3.38, W - M * 2, 0.012, LINE)
for i, (b, a, bb, g, cm) in enumerate([
        ("lululemon", "18", "244", "6 pts", "0% vs 1%"),
        ("New Balance", "4", "246", "6 pts", "0% vs 0%"),
        ("Nike", "8", "227", "10 pts", "0% vs 3%"),
        ("Under Armour", "7", "153", "10 pts", "0% vs 4%")]):
    y = 3.52 + i * 0.52
    for j, val in enumerate([b, a, bb, g, cm]):
        text(s, val, xs[j], y, 2.1, 0.35, size=13,
             bold=(j == 0 or j == 4), color=SIG if j == 4 else INK)
    rect(s, M, y + 0.4, W - M * 2, 0.008, LINE)
rect(s, M, 5.75, 5.6, 1.05, TINT)
text(s, "0 to 4 points", M + 0.3, 5.92, 5.0, 0.4, size=20, font=HEAD, bold=True, color=SIG)
text(s, "agreement on community, all four brands", M + 0.3, 6.35, 5.0, 0.3, size=12, color=SIG)
text(s, "Different collector. Different method. Different brands.\nDifferent years. Same answer.",
     7.0, 5.85, 5.4, 0.9, size=15, bold=True, color=INK, spacing=22)
notes(s, "This is the answer to 'is that just your sample?'. Two independent collections agree. "
         "Do not claim a trend from the second set: it is year-stratified, so its brand mix "
         "swings between years.")

# ── 6 · THE PROOF GAP ───────────────────────────────────────────────────────
s = slide()
eyebrow(s, "The second finding")
title(s, "Two in three ads make a claim and prove nothing")
rect(s, M, 2.3, 5.5, 2.95, TINT)
text(s, "67%", M + 0.4, 2.55, 4.7, 1.4, size=72, font=HEAD, bold=True, color=SIG)
text(s, "of ads carry no hard evidence: no rating, review count, price or guarantee",
     M + 0.4, 3.95, 4.7, 1.0, size=15, color=INK)
text(s, "Where proof does appear, it is wildly uneven", 7.0, 2.3, 5.6, 0.4, size=15, bold=True)
for i, (b, v) in enumerate([("Alo Yoga", "21,615 reviews"), ("Vuori", "13,593 reviews"),
                            ("Nike", "6,660 reviews"), ("Under Armour", "305 reviews"),
                            ("Arc’teryx", "24 reviews"), ("Tracksmith", "22 reviews")]):
    y = 2.85 + i * 0.4
    hot = i > 3
    text(s, b, 7.0, y, 2.7, 0.32, size=13, color=SIG if hot else INK, bold=hot)
    text(s, v, 9.6, y, 3.0, 0.32, size=13, color=SIG if hot else MID, bold=hot,
         align=PP_ALIGN.RIGHT)
    rect(s, 7.0, y + 0.33, 5.6, 0.008, LINE)
text(s, "A premium brand advertising on 22 reviews is not using proof. It is exposing itself.",
     M, 5.6, W - M * 2, 0.4, size=15, bold=True)
notes(s, "67% carry no hard evidence after the audited proof cleanup. Also in the corrected data: "
         "Under Armour runs price language in 86% of its ads. Ratings are the category default proof.")

# ── 7 · THE OPENING ─────────────────────────────────────────────────────────
s = slide()
eyebrow(s, "The opening")
title(s, "One brand owns belonging. Nobody else contests it.")
text(s, "Community is uncontested territory, and even its owner does not prove it.", M, 2.12, W - M * 2, 0.45,
     size=17, bold=True, color=SIG)
for i, (h, b) in enumerate([
        ("What exists", "Gymshark has made belonging half its message: coached couch-to-5K and "
                        "half-marathon programmes. Alo has ALO Access. Everyone else: zero."),
        ("What is missing", "Neither publishes a participation number, a completion rate or a "
                            "member count. The claim exists; the proof does not."),
        ("The move", "A named programme with published numbers, placed where competitors put "
                     "discount percentages. Ten brands have left this lane to one competitor "
                     "who is not defending it with evidence.")]):
    x = M + i * 4.02
    hot = i == 2
    rect(s, x, 2.82, 3.75, 2.8, TINT if hot else MIST)
    text(s, h, x + 0.3, 3.08, 3.15, 0.4, size=15, bold=True, color=SIG if hot else INK)
    text(s, b, x + 0.3, 3.58, 3.15, 1.9, size=12.5, color=SIG if hot else MID)
text(s, "Spend against runners and gym training, not the everyday athleisure block where "
        "lululemon, Vuori and Fabletics are stacked.",
     M, 5.95, W - M * 2, 0.5, size=14, color=MID)
notes(s, "Hostile question: has Gymshark not already won this? Answer: it owns the claim, not the "
         "proof — no numbers published, framed narrowly around the weight room, and ten brands ceded the lane.")

# ── 8 · WHAT WE GOT WRONG ───────────────────────────────────────────────────
s = slide(WASH)
eyebrow(s, "What we got wrong")
title(s, "Our first answer was confident and false")
text(s, "At 42 ads we concluded three brands cited no proof at all. That was wrong.",
     M, 2.1, W - M * 2, 0.45, size=17, color=SIG)
for lbl, x, w in [("Brand", M, 2.2), ("What we said", M + 2.3, 4.4),
                  ("What they actually run", M + 7.0, 4.6)]:
    text(s, lbl.upper(), x, 2.88, w, 0.3, size=10, bold=True, color=GREY)
for i, (b, said, real) in enumerate([
        ("Vuori", "“The biggest proof gap in the group”",
         "4.6 from 13,593 reviews, 100% product guarantee"),
        ("Gymshark", "“Free returns only”", "4.6 from 200 reviews"),
        ("lululemon", "“Fabric adjectives only”", "Three separate ratings")]):
    y = 3.32 + i * 0.86
    text(s, b, M, y, 2.2, 0.45, size=14, bold=True, color=INK)
    text(s, said, M + 2.3, y, 4.4, 0.6, size=12.5, color=GREY, italic=True)
    text(s, real, M + 7.0, y, 4.6, 0.7, size=12.5, color=SIG, bold=True)
    rect(s, M, y + 0.7, W - M * 2, 0.008, SOFTB)
text(s, "The Transparency Center shows four ads until you click “See all ads.” We went back and "
        "collected the full lists. Then a teammate audited every tag and found 42 more errors, "
        "including in our own headline count. All corrected before this presentation.",
     M, 6.15, W - M * 2, 0.7, size=13, color=MID)
notes(s, "Do not skip this. Showing you caught your own false finding is the most credible thing "
         "in the deck and it inoculates you against anyone who checks the source.")

# ── 9 · WHAT WE LEARNED ─────────────────────────────────────────────────────
s = slide()
eyebrow(s, "What we learned")
title(s, "Five things the workbook does not show")
for i, (h, b) in enumerate([
        ("The bottleneck was collection, not intelligence",
         "Classifying 101 ads took ninety seconds. Collecting them took hours."),
        ("A shallow sample gives a confident wrong answer",
         "Not a vague one. That is far more dangerous."),
        ("Model output drifts unless you freeze it",
         "Two runs on identical data moved 14 tags and rewrote the headline."),
        ("Most search advertising is not positioning",
         "Only 16 of 101 ads were brand-level. The rest is product feed."),
        ("The valuable output was the empty column",
         "Every recommendation came from what nobody said.")]):
    y = 2.18 + i * 0.92
    text(s, f"{i+1:02d}", M, y, 0.6, 0.35, size=16, font=HEAD, bold=True, color=SIG)
    text(s, h, M + 0.78, y - 0.03, 6.2, 0.42, size=15, bold=True)
    text(s, b, M + 7.1, y - 0.01, 4.9, 0.62, size=12.5, color=MID)
notes(s, "Close the analytical section. These transfer to any competitive research, "
         "which is the point for the class.")

# ── 10 · CLOSE ──────────────────────────────────────────────────────────────
s = slide(WASH)
text(s, "Everyone is selling the same fabric.\nNobody is selling belonging with proof.",
     M, 2.0, 11.6, 2.1, size=38, font=HEAD, bold=True, color=INK, spacing=46)
rect(s, M, 4.25, 1.5, 0.045, SIG)
text(s, "That is the lane.", M, 4.55, 8, 0.5, size=20, color=SIG, bold=True)
text(s, "Workbook, data, tagging and full write-up:", M, 5.6, 9, 0.3, size=12, color=GREY)
text(s, "github.com/spdooley5949/Creative-Teardown-Marketing", M, 5.92, 9.5, 0.35,
     size=14, bold=True, color=INK)
text(s, "Shane Dooley   ·   with Andrew Silver, Catherine Chambers, Jessica Smith, "
        "Jesse Padilla and William Levine",
     M, 6.62, 11.6, 0.3, size=11, color=GREY)
notes(s, "Land on the lane, then take questions. Ready answers: how do you know the tags are "
         "right; why only Google search; why is Patagonia only three ads; isn't community just "
         "Gymshark's brand voice.")

prs.save(sys.argv[1])
print("wrote", sys.argv[1], "|", len(prs.slides._sldIdLst), "slides")
